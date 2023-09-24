import dash
import pandas as pd
import plotly.graph_objs as go
import plotly.express as px
import plotly.io as pio # Importa el tema oscuro de Plotly
import pytz
import dash_leaflet as dl
import dash_leaflet.express as dlx
from dash import dcc, html, Input, Output
from datetime import datetime
import calendar
import pyautogui
import io
import base64
from PIL import Image
import openpyxl
import img2pdf
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import webbrowser
import dash_table
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt
from utils import (
    get_database_connection,
    obtener_numero_de_pozos,
    generate_data_card,
    get_last_updated_time,
)

external_stylesheets = [
    'https://fonts.googleapis.com/css?family=Lato',
    '/assets/style.css',
]
app = dash.Dash(__name__, external_stylesheets=external_stylesheets)
server = app.server

# ---------------------------------------------------------------------------------------------------------------------------------------------------
# SECCION DE CONEXIONES Y CONSULTAS MYSQL
# ---------------------------------------------------------------------------------------------------------------------------------------------------

# Conexión con el servidor MySQL
conexion = get_database_connection()

# Consulta SQL para obtener los datos sumados por mes
query = """SELECT 
    YEAR(Date) AS Año,
    MONTH(Date) AS Mes,
    AVG(Presion_intake) AS Promedio_Presion_intake,
    AVG(Freq) AS Promedio_Freq,
    AVG(Caudal) AS Promedio_Caudal,
    SUM(WOR) AS Suma_WOR,
    SUM(WCUT) AS Suma_WCUT,
    SUM(BWPD) AS Suma_BWPD,
    SUM(BOPD) AS Suma_BOPD
FROM critical_variables_updated
GROUP BY Año, Mes
"""

# Leer los datos desde MySQL y crear un DataFrame
df = pd.read_sql(query, conexion)

queryrunlife = """
SELECT 
    cvu.Well_Id,
    cvu.Date,
    MAX(cvu.RunLife) AS Maximo_Antes_de_1_0,
    wmu.UWI,
    wmu.Sistema_Levantamiento
FROM (
    SELECT *,
        SUM(CASE WHEN RunLife = 1.0 THEN 1 ELSE 0 END) OVER (PARTITION BY Well_Id ORDER BY Date) AS Cycle
    FROM critical_variables_updated
    WHERE RunLife > 1.0
        AND DATE_FORMAT(Date, '%H:%i') = '00:00'
        AND (SELECT MAX(RunLife) FROM critical_variables_updated WHERE Well_Id = critical_variables_updated.Well_Id AND RunLife = 1.0) IS NOT NULL
) AS cvu
JOIN wells_master_updated AS wmu ON cvu.Well_Id = wmu.Well_Id
GROUP BY cvu.Well_Id, cvu.Date, wmu.UWI, wmu.Sistema_Levantamiento;
"""

# Ejecuta la consulta y guarda los resultados en un DataFrame de pandas
df_runlife = pd.read_sql_query(queryrunlife, conexion)

queryrunstatus = """
    SELECT cvu.Well_Id, DATE_FORMAT(cvu.Date, '%Y-%m-%d') AS Date, cvu.Run_Status, wm.UWI, wm.Sistema_Levantamiento
    FROM critical_variables_updated cvu
    INNER JOIN (
        SELECT Well_Id, MAX(Date) AS MaxDate
        FROM critical_variables_updated
        GROUP BY Well_Id
    ) max_dates
    ON cvu.Well_Id = max_dates.Well_Id AND cvu.Date = max_dates.MaxDate
    INNER JOIN wells_master_updated wm
    ON cvu.Well_Id = wm.Well_Id;
    """

# Suponiendo que ya tienes una conexión a tu base de datos llamada 'conexion'
df_runstatus = pd.read_sql_query(queryrunstatus, conexion)

querypruebas = """
    SELECT  wm.UWI, dpp.Test_Date, dpp.Test_Num, dpp.Certified, dpp.Test_Duration, dpp.API, dpp.BSW_P, wm.Sistema_Levantamiento
    FROM data_prueba_pozo_updated dpp
    INNER JOIN (
        SELECT Well_Id, MAX(Test_Date) AS MaxDate
        FROM data_prueba_pozo_updated
        GROUP BY Well_Id
    ) max_dates
    ON dpp.Well_Id = max_dates.Well_Id AND dpp.Test_Date = max_dates.MaxDate
    INNER JOIN wells_master_updated wm
    ON dpp.Well_Id = wm.Well_Id;
    """

# Suponiendo que ya tienes una conexión a tu base de datos llamada 'conexion'
df_pruebas = pd.read_sql_query(querypruebas, conexion)

# Consulta SQL para obtener los datos

query_var_3 = """
SELECT a.Well_Id,
	   YEAR(Volume_Date) AS Año,
       MONTH(Volume_Date) AS Mes,
       SUM(Oil) AS Suma_Oil,
       SUM(Gas) AS Suma_Gas,
       AVG(`OIL QUALITY`) AS PROMEDIO_OIL_QUALITY,
       b.UWI
FROM data_diaria_volumetrica_updated a
JOIN wells_master_updated b
ON a. Well_Id = b. Well_Id
GROUP BY UWI, Mes, Año
"""

query_tarjetas_var_2 = "SELECT Caudal, WOR FROM critical_variables_updated"
query_tarjetas_var_3 = "SELECT Hours, Oil, Water, Gas FROM data_diaria_volumetrica_updated"

querymap = "SELECT UWI, Geo_latitude, Geo_longitude, Wellhead_depth, Water_depth FROM wells_master_updated"

# Obtener datos
df_query_var_3 = pd.read_sql(query_var_3, conexion)
df_var_2 = pd.read_sql(query_tarjetas_var_2, conexion)
df_var_3 = pd.read_sql(query_tarjetas_var_3, conexion)

# Obtener datos de producción de Oil y Gas por cada Well_Id
well_production_data = df_query_var_3

df_map = pd.read_sql(querymap, conexion)

# Obtener la hora actual en Colombia
colombia_tz = pytz.timezone('America/Bogota')
formatted_time = get_last_updated_time(colombia_tz)

# Obtener el número de pozos
numero_de_pozos = obtener_numero_de_pozos(conexion)

# Consulta SQL para obtener las sumas agrupadas por "Well Id"
query_sumas_por_pozo = """
SELECT a.Well_Id,
       SUM(Gas) AS Total_Gas_Pozo,
       SUM(Oil) AS Total_Oil_Pozo,
       SUM(Water) AS Total_Water_Pozo,
       b.UWI,
       b.Geo_latitude,
       b.Geo_longitude,
       b.Sistema_Levantamiento,
       b.Purpose
FROM data_diaria_volumetrica_updated a
JOIN wells_master_updated b
ON a. Well_Id = b. Well_Id
GROUP BY Well_Id
"""

# Ejecutar la consulta SQL y cargar los resultados en un DataFrame
df_sumas_por_pozo = pd.read_sql(query_sumas_por_pozo, conexion)

# Calcular las proporciones por pozo
df_sumas_por_pozo['Proporcion_Gas_Pozo'] = (df_sumas_por_pozo['Total_Gas_Pozo'] / (df_sumas_por_pozo['Total_Gas_Pozo'] + df_sumas_por_pozo['Total_Oil_Pozo'] + df_sumas_por_pozo['Total_Water_Pozo'])) * 100
df_sumas_por_pozo['Proporcion_Oil_Pozo'] = (df_sumas_por_pozo['Total_Oil_Pozo'] / (df_sumas_por_pozo['Total_Gas_Pozo'] + df_sumas_por_pozo['Total_Oil_Pozo'] + df_sumas_por_pozo['Total_Water_Pozo'])) * 100
df_sumas_por_pozo['Proporcion_Water_Pozo'] = (df_sumas_por_pozo['Total_Water_Pozo'] / (df_sumas_por_pozo['Total_Gas_Pozo'] + df_sumas_por_pozo['Total_Oil_Pozo'] + df_sumas_por_pozo['Total_Water_Pozo'])) * 100

# Consulta SQL
query_hours = """
SELECT
    b.UWI,
    SUM(a.Hours) AS Total_Hours,
    COUNT(DISTINCT a.VOLUME_DATE) AS Days_with_Values,
    SUM(a.Hours) / COUNT(DISTINCT a.VOLUME_DATE) AS Average_Hours_Per_Day
FROM
    tablero_geohallitians.data_diaria_volumetrica_updated AS a
INNER JOIN 
    wells_master_updated AS b ON a.Well_Id = b.Well_Id
WHERE
    a.Hours IS NOT NULL
GROUP BY
    a.Well_Id, b.UWI;
"""

# Cargar datos en un DataFrame de Pandas
df_hours = pd.read_sql(query_hours, conexion)

def generate_pie_chart(row, text_size=11):
    labels = ['Gas', 'Oil', 'Water']
    values = [row['Total_Gas_Pozo'], row['Total_Oil_Pozo'], row['Total_Water_Pozo']]
    colors = ['#663366', '#FF6633', '#0099FF']  # Colores personalizados: morado, naranja y azul

    fig = go.Figure(data=[go.Pie(
        labels=labels,
        values=values,
        hole=0.3,
        marker=dict(colors=colors),
        textinfo='label+percent',  # Muestra etiquetas y porcentaje
        textfont=dict(size=text_size, family='Arial, sans-serif', color='black'),  # Configura la fuente del texto
    )])

    fig.update_layout(
        showlegend=False,  # Oculta la leyenda
        margin=dict(l=0, r=0, t=0, b=0),  # Ajusta los márgenes
        width=150,  # Ajusta el ancho del gráfico
        height=150,  # Ajusta la altura del gráfico
    )

    return dcc.Graph(figure=fig, config={'displayModeBar': False})

# Consulta SQL
query_prom_bopd = """
SELECT
    a.Well_Id,
    AVG(a.BOPD) AS Average_BOPD,
    b.UWI
FROM
    tablero_geohallitians.critical_variables_updated AS a
INNER JOIN 
    wells_master_updated AS b ON a.Well_Id = b.Well_Id
GROUP BY
    a.Well_Id, b.UWI;
"""

# Cargar datos en un DataFrame de Pandas
df_prom = pd.read_sql(query_prom_bopd, conexion)

query_bopd = """
SELECT Well_Id, BOPD, Day
FROM critical_variables_updated
GROUP BY Well_Id, BOPD, Day;
"""

df_bopd = pd.read_sql(query_bopd, conexion)

# Realiza la consulta SQL
query_minus_gas = """
SELECT
    a.Well_Id,
    AVG(a.Gas) AS Average_Gas,
    b.UWI
FROM
    tablero_geohallitians.data_diaria_volumetrica_updated AS a
INNER JOIN 
    wells_master_updated AS b ON a.Well_Id = b.Well_Id
GROUP BY
    a.Well_Id, b.UWI;
"""
df_minus_gas = pd.read_sql_query(query_minus_gas, conexion)

# Realiza la consulta SQL
query_minus_oil = """
SELECT
    a.Well_Id,
    AVG(a.BOPD) AS Average_BOPD,
    b.UWI
FROM
    tablero_geohallitians.critical_variables_updated AS a
INNER JOIN 
    wells_master_updated AS b ON a.Well_Id = b.Well_Id
GROUP BY
    a.Well_Id, b.UWI;
"""
df_minus_oil = pd.read_sql_query(query_minus_oil, conexion)


# Realiza la consulta SQL con filtrado por UWI
query_production_bopd = """
SELECT
    a.Well_Id,
    YEAR(Date) AS Año,
    MONTH(Date) AS Mes,
    SUM(a.BOPD) AS Suma_BOPD,
    b.UWI
FROM
    tablero_geohallitians.critical_variables_updated AS a
INNER JOIN 
    wells_master_updated AS b ON a.Well_Id = b.Well_Id
WHERE
    b.UWI IN ('Well013', 'Well017', 'Well002', 'Well008', 'Well010')
GROUP BY
    a.Well_Id, b.UWI, Mes, Año;
"""

well_production_bopd = pd.read_sql_query(query_production_bopd, conexion)

# ---------------------------------------------------------------------------------------------------------------------------------------------------
# SECCION DE GRAFICAS Y TARJETAS
# ---------------------------------------------------------------------------------------------------------------------------------------------------

# Estilos comunes de las gráficas de dispersión
common_style = {
    'display': 'inline-block',
    'backgroundColor': '#000000',
    'width': '45%',
    'height': '350px',
    'color': 'white',
}

# Ordena el DataFrame por Average_Gas de manera ascendente y toma las primeras 5 filas
df_minus_gas = df_minus_gas.sort_values(by='Average_Gas', ascending=True).head(5)

def create_bar_chart(df_minus_gas):
    # Ordena el DataFrame por Average_Gas de manera ascendente y toma las primeras 5 filas
    df_minus_gas = df_minus_gas.sort_values(by='Average_Gas', ascending=True).head(5)

    # Crea el gráfico de barras invertidas con márgenes personalizados
    fig = go.Figure(data=[go.Bar(
        x=df_minus_gas['Average_Gas'],
        y=df_minus_gas['UWI'],
        orientation='h',
    )])
    
    # Establece el color de las barras
    fig.update_traces(marker=dict(color='rgba(0, 204, 255, 0.2)'))
    
    # Aplica márgenes personalizados al gráfico y ajusta el espaciado del eje Y
    fig.update_layout(
        margin=dict(l=10, r=10, t=30, b=10),
        plot_bgcolor='#000000',
        paper_bgcolor='#000000',
        xaxis_title=dict(text='Average_Gas', font=dict(color='white')),  # Cambia el color del texto del eje X
        yaxis_title=dict(text='UWI', font=dict(color='white')),  # Cambia el color del texto del eje Y
        xaxis_tickfont=dict(color='white'),  # Cambia el color de las etiquetas de los ticks del eje X
        yaxis_tickfont=dict(color='white'),  # Cambia el color de las etiquetas de los ticks del eje Y
        yaxis_fixedrange=True,  # Establece el rango fijo del eje Y para ajustar el espaciado
    )
    
    return dcc.Graph(figure=fig, id='bar-chart', style={'maxHeight': '235px', 'overflowY': 'scroll'})



# Agrega el título de la tarjeta aquí
top5minusgas_card = html.Div(
    style={'height': '300px', 'textAlign': 'center', 'align-content': 'right', 'align-items': 'right', 'border-radius':'7px', 'box-shadow': '0px 0px 5px 2px rgba(255, 255, 255, 0.2)', 'backgroundColor': '#131313', 'margin-bottom': '30px'},
    children=[
        html.Div(
            "Top 5 lowest gas production",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        create_bar_chart(df_minus_gas)  # Llama a la función para crear el gráfico
    ],
)

# Ordena el DataFrame por Average_Gas de manera ascendente y toma las primeras 5 filas
df_minus_oil = df_minus_oil.sort_values(by='Average_BOPD', ascending=True).head(5)

def create_bar_oil(df_minus_oil):
    # Ordena el DataFrame por Average_BOPD de manera ascendente y toma las primeras 5 filas
    df_minus_oil = df_minus_oil.sort_values(by='Average_BOPD', ascending=True).head(5)
    
    # Crea el gráfico de barras invertidas con márgenes personalizados
    fig = go.Figure(data=[go.Bar(
        x=df_minus_oil['Average_BOPD'],
        y=df_minus_oil['UWI'],
        orientation='h',
    )])
    
    # Establece el color de las barras
    fig.update_traces(marker=dict(color='rgba(0, 204, 255, 0.2)'))
    
    # Aplica márgenes personalizados al gráfico y ajusta el espaciado del eje Y
    fig.update_layout(
        margin=dict(l=10, r=10, t=30, b=10),
        plot_bgcolor='#000000',
        paper_bgcolor='#000000',
        xaxis_title=dict(text='Average Oil', font=dict(color='white')),  # Cambia el color del texto del eje X
        yaxis_title=dict(text='UWI', font=dict(color='white')),  # Cambia el color del texto del eje Y
        xaxis_tickfont=dict(color='white'),  # Cambia el color de las etiquetas de los ticks del eje X
        yaxis_tickfont=dict(color='white'),  # Cambia el color de las etiquetas de los ticks del eje Y
        yaxis_fixedrange=True,  # Establece el rango fijo del eje Y para ajustar el espaciado
    )
    
    return dcc.Graph(figure=fig, id='bar-charto', style={'maxHeight': '235px', 'overflowY': 'scroll'})



# Agrega el título de la tarjeta aquí
top5minusoil_card = html.Div(
    style={'height': '300px', 'textAlign': 'center', 'align-content': 'right', 'align-items': 'right', 'border-radius':'7px', 'box-shadow': '0px 0px 5px 2px rgba(255, 255, 255, 0.2)', 'backgroundColor': '#131313', 'margin-bottom': '30px'},
    children=[
        html.Div(
            "Top 5 lowest oil production",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        create_bar_oil(df_minus_oil)  # Llama a la función para crear el gráfico
    ],
)

# Lista de nombres de tablas
tablas = ["data_prueba_pozo_updated", "wells_master_updated", "data_diaria_volumetrica_updated", "critical_variables_updated"]

def generate_heatmap(tables):
    cursor = conexion.cursor()

    resultado_final = pd.DataFrame(columns=["Tabla", "Columna", "Total nulos"])
    data_traces = []  # Lista para almacenar objetos de datos (traces)

    for tabla in tables:
        # Obtiene la lista de columnas en la tabla
        cursor.execute(f"DESCRIBE {tabla}")
        columnas = [columna[0] for columna in cursor.fetchall()]

        # Crea una consulta SQL dinámica para contar valores nulos en cada columna
        consulta = "SELECT COUNT(*) AS TotalRegistros, " + ", ".join([
            f"SUM(CASE WHEN `{col}` IS NULL THEN 1 ELSE 0 END) AS `{col}_Nulos`"
            for col in columnas
        ]) + f" FROM {tabla}"

        # Ejecuta la consulta
        cursor.execute(consulta)

        # Obtiene los resultados
        resultado = cursor.fetchone()

        # Crear un DataFrame a partir de los resultados
        data = []
        for col, value in zip(columnas, resultado[1:]):
            total_vacios = value
            data.append([tabla, col, int(total_vacios)])

        df = pd.DataFrame(data, columns=["Tabla", "Columna", "Total nulos"])

        # Agregar los resultados al DataFrame final
        resultado_final = pd.concat([resultado_final, df], ignore_index=True)

        # Convertir la columna "Total nulos" a tipo numérico
        resultado_final["Total nulos"] = pd.to_numeric(resultado_final["Total nulos"])

        # Modificar nombres de columnas (eliminar "_updated")
        resultado_final["Columna"] = resultado_final["Columna"].str.replace("_updated", "")

        # Crear un mapa de calor con paleta cálida (YlOrRd) usando Plotly Express
        fig = px.imshow(resultado_final.pivot_table(index="Tabla", columns="Columna", values="Total nulos", aggfunc="sum"),
                        labels=dict(x="Columnas", y="Tablas", color="Total nulos"), color_continuous_scale="YlOrRd")
        
        # Agregar el objeto de datos (trace) a la lista
        data_traces.append(fig['data'][0])
        
    return data_traces  # Devolver la lista de objetos de datos (traces)

# Agrega el título de la tarjeta aquí
heatmap_card = html.Div(
    style={'height': '100%', 'width': '100%', 'textAlign': 'center', 'align-content': 'right', 'align-items': 'right', 'border-radius':'7px', 'box-shadow': '0px 0px 5px 2px rgba(255, 255, 255, 0.2)', 'backgroundColor': '#131313', 'margin-bottom': '30px'},
    children=[
        html.Div(
            "Missing Values",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        dcc.Graph(
            id='heatmap-graph',
            figure={
                'data': generate_heatmap(tablas),  # Llama a tu función para generar el heatmap
                'layout': go.Layout(
                    plot_bgcolor='black',  # Cambia el fondo del gráfico a negro
                    paper_bgcolor='black',  # Cambia el fondo del área de papel a negro
                    margin={'t': 30, 'l': 100, 'r': 50, 'b': 70},
                    coloraxis_colorbar=dict(title='Total Nulls', title_font=dict(color='white')),  # Color del texto de la leyenda
                    xaxis=dict(tickangle=45, tickfont=dict(size=9)),
                    yaxis=dict(tickangle=45, tickfont=dict(size=8))
                ),
            },
            style={'height': '400px', 'width': '1250px'}  # Define el alto y ancho deseados para la gráfica
        )
    ],
)


# Unir los dataframes en función del campo "Well_Id"
merged_df = df_bopd.merge(df_prom, on='Well_Id', how='inner')

# Calcular el umbral del 10% del valor PROMEDIO_BOPD
merged_df['Umbral'] = 0.10 * merged_df['Average_BOPD']

# Filtrar los valores de BOPD que son menores que el umbral
filtered_df = merged_df[merged_df['BOPD'] < merged_df['Umbral']]

# Estilo para la tabla
table_style = {
    'width': '240px',
    'border-collapse': 'collapse',
}

# Estilo para las celdas de la tabla, incluyendo el tamaño de fuente
cell_style = {
    'padding': '4px',
    'text-align': 'left',
    'color': 'white',
    'font-size': '9px',  # Tamaño de fuente deseado, puedes ajustarlo según tus preferencias
    'height': 'auto',  # Altura automática basada en el contenido
}

# Estilo para las celdas de encabezado de la tabla, incluyendo el tamaño de fuente
header_cell_style = {
    'padding': '2px',
    'text-align': 'left',
    'color': 'black',
    'font-size': '9px',  # Tamaño de fuente deseado para las celdas de encabezado
}

# Estilo para el contenedor de la tabla
table_container_style = {
    'overflowY': 'auto',  # Agrega scroll vertical si es necesario
    'maxHeight': '270px',  # Establece la altura máxima de la tabla
    'margin': '0px',  # Margen cero
}

# Definir el diseño de la tabla en la función tabla_critical
def tabla_critical():
    # Filtra el DataFrame para incluir solo las columnas deseadas y en el orden deseado
    filtered_df_subset = filtered_df[['UWI', 'BOPD', 'Average_BOPD', 'Day']]
    
    # Redondear la columna "Average BOPD" a 2 decimales
    filtered_df_subset['Average_BOPD'] = filtered_df_subset['Average_BOPD'].round(2)
    
    # Ordena el DataFrame por la columna "Day" de mayor a menor
    filtered_df_subset = filtered_df_subset.sort_values(by='Day', ascending=False)
    
    # Agrega una nueva columna "Type" con el valor "bopd production less than 10% average"
    filtered_df_subset['Type'] = "bopd production less than 10% average"
    
    return html.Div(
        style=table_container_style,  # Aplicar estilo al contenedor de la tabla
        children=[
            html.Table(
                # Encabezados de la tabla
                [html.Tr([html.Th(col, style=header_cell_style) for col in filtered_df_subset.columns])],
                style=table_style,
            )] +
            # Filas de datos
            [html.Table(
                # Celdas de datos con estilo de fuente
                [html.Tr([html.Td(filtered_df_subset.iloc[i][col], style=cell_style) for col in filtered_df_subset.columns])],
                style=table_style,
            ) for i in range(len(filtered_df_subset))]
    )

# Agrega el título de la tarjeta aquí
alertas_card = html.Div(
    style={'height': '300px', 'textAlign': 'center', 'align-content': 'right', 'align-items': 'right', 'border-radius':'7px', 'box-shadow': '0px 0px 5px 2px rgba(255, 255, 255, 0.2)', 'backgroundColor': 'rgba(255, 0, 0, 0.4)', 'margin-bottom': '30px'},
    children=[
        html.Div(
            "Critical Alerts",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        tabla_critical()
    ],
)


# Calcular el promedio de "Average_Hours_Per_Day"
average_hours_per_day = df_hours['Average_Hours_Per_Day'].mean()

def hours():
    # Crear la figura del velocímetro
    figure = {
        'data': [
            go.Indicator(
                mode="gauge+number",
                value=average_hours_per_day,
                domain={'x': [0, 1], 'y': [0, 1]},
                title={'text': "Average Hours Per Day", 'font': {'color': 'white'}},  # Estilo para el título
                gauge={
                    'axis': {
                        'range': [None, max(df_hours['Average_Hours_Per_Day'])],  # Cambio aquí
                        'tickvals': [0, 5, 10, 15, 20],
                        'tickfont': {'color': 'white'}
                    },
                    'steps': [
                        {'range': [0, max(df_hours['Average_Hours_Per_Day'])], 'color': "white"},  # Cambio aquí
                        {'range': [0, average_hours_per_day], 'color': "green"}
                    ],
                },
                number={'font': {'color': 'white'}}
            )
        ],
        'layout': {
            'plot_bgcolor': 'black',
            'paper_bgcolor': 'black',
            'margin': {'t': 25, 'l': 10, 'r': 70, 'b': 7}
        },
    }

    return dcc.Graph(id='velocimeter', style={'width': '320px', 'height': '110px','overflowY': 'auto', 'overflow': 'hidden', 'overflowX': 'auto'}, figure=figure)


# Agrega el título de la tarjeta aquí
hours_card = html.Div(
    style={'height': '150px', 'textAlign': 'center', 'align-content': 'right', 'align-items': 'right', 'border-radius':'7px', 'box-shadow': '0px 0px 5px 2px rgba(255, 255, 255, 0.2)', 'backgroundColor': '#131313', 'margin-bottom': '30px'},
    children=[
        html.Div(
            "Average Hours/Day",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        hours(),
    ],
)

def tabla_runlife():

    # Identificar los índices donde Maximo_Antes_de_1_0 es 2.0 (inicio de un nuevo ciclo)
    new_cycle_indices = df_runlife[df_runlife['Maximo_Antes_de_1_0'] == 2.0].index

    # Crear un array con los valores de ciclo actual
    cycle_values = np.zeros(len(df_runlife))

    current_cycle = 1  # Inicializar el contador de ciclos

    # Modificar los valores de ciclo para que comiencen un nuevo ciclo
    for index in new_cycle_indices:
        cycle_values[index:] = current_cycle
        current_cycle += 1  # Incrementar el contador de ciclos

    # Agregar la columna 'Cycle' al DataFrame
    df_runlife['Cycle'] = cycle_values

    # Calcular los máximos valores por ciclo para cada pozo
    max_values_by_cycle = df_runlife.groupby(['Well_Id', 'Cycle'])['Maximo_Antes_de_1_0'].max()

    # Crear un DataFrame con los máximos valores por ciclo
    max_values_df = max_values_by_cycle.reset_index()

    # Agregar la columna 'Date' al DataFrame max_values_df
    max_values_df['Date'] = df_runlife.groupby(['Well_Id', 'Cycle'])['Date'].max().values

    # Agregar la columna 'Date' al DataFrame max_values_df y formatearla en el formato "año-mes-día"
    max_values_df['Date'] = pd.to_datetime(max_values_df['Date']).apply(lambda x: x.strftime('%Y-%m-%d'))

    # Obtener las columnas 'UWI' y 'Sistema_Levantamiento' correspondientes al valor máximo de 'Maximo_Antes_de_1_0' en cada ciclo
    uwis_sistema_levantamiento = df_runlife[df_runlife.groupby(['Well_Id', 'Cycle'])['Maximo_Antes_de_1_0'].transform('max') == df_runlife['Maximo_Antes_de_1_0']][['Well_Id', 'Cycle', 'UWI', 'Sistema_Levantamiento']]

    # Fusionar los DataFrames para agregar 'UWI' y 'Sistema_Levantamiento' a max_values_df
    max_values_df = max_values_df.merge(uwis_sistema_levantamiento, on=['Well_Id', 'Cycle'])

    # Eliminar la columna 'Cycle'
    max_values_df = max_values_df.drop(columns=['Cycle'])

    # Cambiar el nombre de la columna 'Maximo_Antes_de_1_0' a 'Run Life'
    max_values_df = max_values_df.rename(columns={'Maximo_Antes_de_1_0': 'Run Life'})

    # Agregar la columna 'Cycle' con valores incrementales por 'Well_Id'
    max_values_df['Cycle'] = max_values_df.groupby('Well_Id').cumcount() + 1

    # Cambiar el nombre de la columna 'Sistema_Levantamiento' a 'Lift_System'
    max_values_df = max_values_df.rename(columns={'Sistema_Levantamiento': 'Lift System'})

    # Trasladar la columna 'UWI' a la primera columna
    max_values_df = max_values_df[['UWI'] + [col for col in max_values_df.columns if col != 'UWI']]

    # Eliminar la columna 'Well_Id'
    max_values_df = max_values_df.drop(columns=['Well_Id'])

    valores_verdes = [1426, 413, 343, 167, 496, 573, 868]

    # Crea la tabla de Dash
    table = dash_table.DataTable(
        id='tabla_runlife',
        columns=[
            {'name': col, 'id': col} for col in max_values_df.columns
        ],
        data=max_values_df.to_dict('records'),
        style_table={'width': '400px', 'maxHeight': '300px', 'overflowY': 'auto'},
        style_header={
            'textAlign': 'center'  # Centrar los encabezados de las columnas
        },
        style_cell={
            'textAlign': 'left', 'backgroundColor': 'black', 'color': 'white',  # Alinear los valores de las celdas a la izquierda
        },
        css=[
            {
                'selector': '.dash-header',
                'rule': 'background-color: #131313 !important; color: white !important;'
            }
        ],
        style_data_conditional=[
            {
                'if': {
                    'column_id': 'Run Life',
                    'filter_query': '{{Run Life}} = {}'.format(valor),
                },
                'backgroundColor': 'rgba(51, 255, 204, 0.4)',  # Verde pastel con transparencia
                'color': 'black',  # Texto en negro
            } for valor in valores_verdes
        ]
    )

    return table

def tabla_runstatus():
    # Crear la tabla Dash y devolverla como un componente
    return dash_table.DataTable(
        id='table',
        columns=[
            {"name": "UWI", "id": "UWI"},
            {"name": "Run Status", "id": "Run_Status"},
            {"name": "Date", "id": "Date"},
            {"name": "Lift System", "id": "Sistema_Levantamiento"},
        ],
        data=df_runstatus.to_dict('records'),
        style_table={'width': '400px', 'maxHeight': '300px', 'overflowY': 'auto'},
        style_header={
            'textAlign': 'center'  # Centrar los encabezados de las columnas
        },
        style_cell={
            'textAlign': 'left', 'backgroundColor': 'black', 'color': 'white',  # Alinear los valores de las celdas a la izquierda
        },
        css=[
            {
                'selector': '.dash-header',
                'rule': 'background-color: #131313 !important; color: white !important;'
            }
        ],
        style_data_conditional=[
            {
                'if': {
                    'column_id': 'Run_Status',
                    'filter_query': '{Run_Status} > 0',
                },
                'backgroundColor': 'rgba(51, 255, 204, 0.4)',  # Verde pastel brillante con transparencia
                'color': 'black',  # Texto en negro
            },
            {
                'if': {
                    'column_id': 'Run_Status',
                    'filter_query': '{Run_Status} = 0',
                },
                'backgroundColor': 'rgba(255, 0, 0, 0.4)',  # Rojo con transparencia
                'color': 'black',  # Texto en negro
            },
        ]
    )

def tabla_pruebas(df_pruebas):
    return dash_table.DataTable(
        id='table-pruebas',
        columns=[
            {"name": "UWI", "id": "UWI"},
            {"name": "Date", "id": "Test_Date"},
            {"name": "#", "id": "Test_Num"},
            {"name": "Certified", "id": "Certified"},
            {"name": "Duration", "id": "Test_Duration"},
            {"name": "API", "id": "API"},
            {"name": "BSW_P", "id": "BSW_P"},
            {"name": "Lift System", "id": "Sistema_Levantamiento"},
        ],
        data=df_pruebas.to_dict('records'),
        style_table={'width': '400px', 'maxHeight': '300px', 'overflowY': 'auto', 'backgroundColor': 'black'},
        style_header={'textAlign': 'center'},
        style_cell={'textAlign': 'left', 'backgroundColor': 'black', 'color': 'white'},
        css=[
            {
                'selector': '.dash-header',
                'rule': 'background-color: #131313 !important; color: white !important;'
            }
        ],
        style_data_conditional=[
            {
                'if': {
                    'column_id': 'Certified',
                    'filter_query': '{Certified} = 1',
                },
                'backgroundColor': 'rgba(51, 255, 204, 0.4)',  # Verde pastel brillante con transparencia
                'color': 'black',  # Texto en negro
            },
            {
                'if': {
                    'column_id': 'Certified',
                    'filter_query': '{Certified} = 0',
                },
                'backgroundColor': 'rgba(255, 0, 0, 0.4)',  # Rojo con transparencia
                'color': 'black',  # Texto en negro
            },
            {
                'if': {
                    'column_id': 'Certified',
                    'filter_query': '{Certified} = ""',
                },
                'backgroundColor': 'rgba(255, 165, 0, 0.6)',  # Naranja con transparencia
                'color': 'black',  # Texto en negro
            },
        ]
    )

colors_UWIS = {
    'Well001': 'rgba(67, 160, 71, 0.6)',
    'Well002': 'rgba(188, 170, 164, 1)',
    'Well004': 'rgba(0, 204, 255, 0.2)',
    'Well005': 'rgba(183, 28, 28, 0.4)',
    'Well006': 'rgba(255, 87, 34, 0.3)',
    'Well007': 'rgba(186, 104, 200, 0.2)',
    'Well008': 'rgba(255, 255, 157, 0.5)',
    'Well009': 'rgba(204, 153, 255, 0.2)',
    'Well010': 'rgba(102, 255, 204, 0.6)',
    'Well012': 'rgba(204, 255, 204, 1)',
    'Well013': 'rgba(255, 255, 204, 0.2)',
    'Well014': 'rgba(255, 255, 153, 0.2)',
    'Well015': 'rgba(255, 153, 204, 0.6)',
    'Well016': 'rgba(153, 204, 102, 0.4)',
    'Well017': 'rgba(153, 153, 204, 1)',
    'Well018': 'rgba(97, 97, 97, 0.6)',
    'Well019': 'rgba(141, 110, 99, 0.6)'
}

# Agrega el título de la tarjeta aquí
efficiency_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Operational efficiency",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        tabla_runlife(),
    ],
)

# Agrega el título de la tarjeta aquí
runstatus_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Run Status updated",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        tabla_runstatus()
    ],
)

# Agrega el título de la tarjeta aquí
Pruebas_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Latest tests",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        tabla_pruebas(df_pruebas)
    ],
)

def formatar_mes_año(mes, año):
    meses_abreviados = {
        1: 'Jan',
        2: 'Feb',
        3: 'Mar',
        4: 'Apr',
        5: 'May',
        6: 'Jun',
        7: 'Jul',
        8: 'Aug',
        9: 'Sep',
        10: 'Oct',
        11: 'Nov',
        12: 'Dec'
    }
    return f"{meses_abreviados.get(mes, 'N/A')}/{año}"


def graficar_tendencia_oil_quality(df):
    # Ordena el DataFrame por la columna 'Mes' antes de la agrupación
    df = df.sort_values(by='Mes')

    # Crear una figura de tendencia utilizando Plotly
    fig = go.Figure()
    
    for pozo_id, data in df.groupby('UWI'):
        data['Mes'] = data.apply(lambda row: formatar_mes_año(row['Mes'], row['Año']), axis=1)  # Formatea el eje X
        figura_pozo = go.Scatter(
            x=data['Mes'],
            y=data['PROMEDIO_OIL_QUALITY'],
            mode='lines+markers',
            name=pozo_id,
            line=dict(color=colors_UWIS.get(pozo_id, 'rgba(255, 255, 255, 1)')),  # Asigna el color según el UWI
        )
        fig.add_trace(figura_pozo)
    
    # Configurar el diseño de la gráfica
    fig.update_layout(
        xaxis_title='Mes',
        yaxis_title='Average Oil Quality',
        width=400,
        height=300,
        margin=dict(l=0, r=120, b=0, t=30),
        plot_bgcolor='#000000',
        paper_bgcolor='#000000',
        font=dict(color='white', size=9),  # Cambia el color y el tamaño del texto
        legend=dict(
            font=dict(color='white', size=9),  # Cambia el color y el tamaño de la leyenda
            traceorder='normal',  # Establece el orden de los elementos de la leyenda como "normal"
            itemsizing='constant',  # Establece el tamaño de los elementos de la leyenda como "constante"
        ),
        xaxis=dict(title=dict(font=dict(color='white', size=12)), showgrid=True, gridcolor='#131313', showline=True, linewidth=1),  # Cambia el color y el tamaño del título del eje X
        yaxis=dict(title=dict(font=dict(color='white', size=12)), showgrid=True, gridcolor='#131313'),  # Cambia el color y el tamaño del título del eje Y
        showlegend=True,  # Muestra la leyenda
    )
    
    return dcc.Graph(
        figure=fig,
        config={'displayModeBar': True},  # Configura displayModeBar aquí
    )

# Suponiendo que ya tienes tu DataFrame df_query_var_3
grafica_oil_quality = graficar_tendencia_oil_quality(df_query_var_3)
    

# Agrega el título de la tarjeta aquí
API_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "API",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        graficar_tendencia_oil_quality(df_query_var_3),  # Inserta el mapa directamente después del título
    ],
)

# Profundidad de pozos

def trajectory_wells(df_map):
    wellhead_depth = df_map['Wellhead_depth']
    water_depth = df_map['Water_depth']
    geo_longitude = df_map['Geo_longitude']
    geo_latitude = df_map['Geo_latitude']
    uwi_list = df_map['UWI']
    wellhead_depth_invertido = [-depth for depth in wellhead_depth]
    water_depth_invertido = [-depth for depth in water_depth]

    fig = go.Figure()

    colors_UWIS = {
        'Well001': 'rgba(67, 160, 71, 0.6)',
        'Well002': 'rgba(188, 170, 164, 1)',
        'Well004': 'rgba(0, 204, 255, 0.2)',
        'Well005': 'rgba(183, 28, 28, 0.4)',
        'Well006': 'rgba(255, 87, 34, 0.3)',
        'Well007': 'rgba(186, 104, 200, 0.2)',
        'Well008': 'rgba(255, 255, 157, 0.5)',
        'Well009': 'rgba(204, 153, 255, 0.2)',
        'Well010': 'rgba(102, 255, 204, 0.6)',
        'Well012': 'rgba(204, 255, 204, 1)',
        'Well013': 'rgba(255, 255, 204, 0.2)',
        'Well014': 'rgba(255, 255, 153, 0.2)',
        'Well015': 'rgba(255, 153, 204, 0.6)',
        'Well016': 'rgba(153, 204, 102, 0.4)',
        'Well017': 'rgba(153, 153, 204, 1)',
        'Well018': 'rgba(97, 97, 97, 0.6)',
        'Well019': 'rgba(141, 110, 99, 0.6)'
    }

    for i in range(len(wellhead_depth)):
        fig.add_trace(go.Scatter3d(
            x=[geo_longitude[i]],
            y=[geo_latitude[i]],
            z=[0],
            mode='markers+text',
            marker=dict(
                size=5,
                color=colors_UWIS.get(uwi_list[i], 'rgba(0, 0, 0, 0.6)'),
                opacity=1
            ),
            textposition='bottom center',
            showlegend=True,
            name=uwi_list[i]
        ))

    # Después de agregar los trazos de la profundidad del pozo, agrega los trazos de profundidad del agua
    for i in range(len(wellhead_depth)):    
        # Trazo para la profundidad del agua (azul)
        fig.add_trace(go.Scatter3d(
            x=[geo_longitude[i], geo_longitude[i]],
            y=[geo_latitude[i], geo_latitude[i]],
            z=[0, water_depth_invertido[i]],  # Utiliza los valores de Water depth
            mode='lines+text',
            line=dict(
                color='#0066FF',  # Color azul para las líneas de profundidad del agua
                width=2
            ),
            showlegend=False,
            name=uwi_list[i] + ' (Water)'
        ))

        # Trazo para la profundidad del pozo (como antes, invertido)
        fig.add_trace(go.Scatter3d(
            x=[geo_longitude[i], geo_longitude[i]],
            y=[geo_latitude[i], geo_latitude[i]],
            z=[0, wellhead_depth_invertido[i]],
            mode='lines+text',
            line=dict(
                color=colors_UWIS.get(uwi_list[i], 'rgba(0, 0, 0, 0.6)'),
                width=2
            ),
            showlegend=False,
            name=uwi_list[i] + ' (Depth)'
    ))

    fig.update_scenes(
        aspectmode='cube'
    )

    fig.update_layout(
        scene=dict(
            xaxis_title='Long',
            yaxis_title='Lat',
            zaxis_title='Depth'
        ),
        scene_camera=dict(
            up=dict(x=0, y=0, z=1),
            center=dict(x=0, y=0, z=0),
            eye=dict(x=1.5, y=1.5, z=0)
        ),
        template="plotly_dark",
        margin=dict(l=0, r=0, b=0, t=20),
        height=300,
        width=400,
        legend=dict(
            font=dict(size=10),
            x=0.7,  # Ajusta la posición horizontal de la leyenda
            y=0.4,   # Ajusta la posición vertical de la leyenda
        ),
        plot_bgcolor='#000000',
        paper_bgcolor='#000000'
    )

    return fig

# Crea una lista de elementos que deseas incluir en wells_depth_card
wells_depth_children = [
    html.Div(
        "Well depth and aquifer depth",
        style={
            'paddingTop': '10px',
            'textAlign': 'center',
            'align-content': 'center',
            'color': 'white',
            'borderBottom': '1px solid white',
            'paddingBottom': '10px',
            'fontWeight': 'bold',
        },
    ),
    html.Div(
        dcc.Graph(
            figure=trajectory_wells(df_map),
            config={
                'displayModeBar': True  # Asegúrate de que displayModeBar esté configurado como True
            }
        ),
        # style={'maxHeight': '300px', 'overflowY': 'auto'}  # Establece una altura máxima y activa la barra de desplazamiento vertical
    ),
]

# Ahora, crea wells_depth_card y pasa la lista de children
wells_depth_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=wells_depth_children,
)


def create_map():

    # Diccionario de mapeo de sistemas de levantamiento a iconos
    sistema_icon_mapping = {
        'BES': 'assets/ESP.png',
        'BME': 'assets/BM.png',
        'GAS LIFT': 'assets/GASLIFT.png',
        'PCP': 'assets/PCP.png'
    }

    # Crea el mapa y los marcadores
    map_ = dl.Map(
        center=[5.533995902516049, -73.35830183000226],
        zoom=6.49,  # Nivel de zoom
        children=[
            dl.TileLayer(url='https://{s}.basemaps.cartocdn.com/dark_all/{z}/{x}/{y}{r}.png')
        ],
        style={'width': '400px', 'height': '300px'}
    )

    # Agrega marcadores para cada pozo en el DataFrame df_sumas_por_pozo
    for _, row in df_sumas_por_pozo.iterrows():  # Utiliza df_sumas_por_pozo en lugar de df_map
        sistema_levantamiento = row['Sistema_Levantamiento']
        purpose = row['Purpose']

        # Obtiene la ruta del icono basado en el sistema de levantamiento
        icon_url = sistema_icon_mapping.get(sistema_levantamiento, 'assets/icono_pozo6.png')

        custom_icon = {
            'iconUrl': icon_url,  # Usa la ruta del icono correspondiente
            'iconSize': [45, 45],  # Tamaño del icono personalizado (ajusta según tus necesidades)
        }

        # Crea el texto para el tooltip con UWI y Purpose
        tooltip_text = f'*UWI: {row["UWI"]}, *Purpose: {purpose}'

        # Crea el marcador con el icono personalizado y el tooltip modificado
        marker = dl.Marker(
            position=[row['Geo_latitude'], row['Geo_longitude']],
            icon=custom_icon,
            children=[
                dl.Tooltip(tooltip_text, direction='top'),  # Muestra UWI y Purpose en el tooltip
                dl.Popup(html.Div([generate_pie_chart(row)]))
            ]
        )

        # Agrega todos los marcadores al mapa usando la propiedad 'children'
        map_.children.append(marker)

    return map_

# Agrega el título de la tarjeta aquí
wells_location_card = html.Div(
    style={'height': '350px', 'width': '60%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Wells Location",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        create_map(),  # Inserta el mapa directamente después del título
    ],
)

def barrel_price():
    return html.Iframe(
        src="https://www.preciopetroleo.net/productos/tv-brent.html",
        width="400px",
        height="600px",
        style={'transform': 'scale(0.72)'},
    )

# Agrega el título de la tarjeta aquí
Barrel_price_realtime_card = html.Div(
    style={
        'height': '400px',
        'textAlign': 'center',
        'border-radius': '10px',
        'box-shadow': '0px 0px 5px 2px rgba(255, 255, 255, 0.2)',
        'backgroundColor': '#131313',
        'overflow': 'hidden',  # Establece el overflow a hidden para recortar contenido
        'overflowX': 'auto',  # Habilita la barra de desplazamiento horizontal cuando sea necesario
    },
    children=[
        html.Div(
            "Barrel Price Real Time (BRENT/WTI)",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        html.Div(  # Contenedor para el iframe con margen superior
            barrel_price(),
            style={'marginTop': '-75px', 'marginLeft': '-60px', 'marginRight': '-55px', 'marginBottom': '-50px'}  # Ajusta los márgenes según sea necesario
        ),
    ],
)

def presion_caudal_graph(df):
    figure = {
        'data': [
            go.Scatter(
                x=df['Promedio_Presion_intake'],
                y=df['Promedio_Caudal'],
                mode='markers',
                marker=dict(
                    size=8,
                    symbol='circle',
                    color='#0066FF',
                    line=dict(
                        width=0.4,
                        color='#212121'
                    )
                ),
                text=df['Mes'],
            )
        ],
        'layout': go.Layout(
            xaxis={'title': 'Average_Presion_Intake', 'color': 'white', 'titlefont': {'color': 'white', 'size': 12}, 'tickfont': {'color': 'white', 'size': 10}},
            yaxis={'title': 'Average Caudal', 'color': 'white', 'titlefont': {'color': 'white', 'size': 12}, 'tickfont': {'color': 'white', 'size': 10}},
            hovermode='closest',
            plot_bgcolor='#000000',
            paper_bgcolor='#000000',
            margin=dict(l=50, r=0, b=50, t=30),
            height=300,
            width=400,
        )
    }

    return dcc.Graph(
        id='presion-caudal-graph',
        config={'displayModeBar': True},
        style={'display': 'inline-block',
            'backgroundColor': '#000000',
            'color': 'white',},
        figure=figure
    )

def freq_caudal_graph(df):
    figure = {
        'data': [
            go.Scatter(
                x=df['Promedio_Freq'],
                y=df['Promedio_Caudal'],
                mode='markers',
                marker=dict(
                    size=8,
                    symbol='circle',
                    color='#0066FF',
                    line=dict(
                        width=0.4,
                        color='#212121'
                    ),
                ),
                text=df['Mes'],
            )
        ],
        'layout': go.Layout(
            xaxis={'title': 'Average Freq', 'color': 'white', 'titlefont': {'color': 'white', 'size': 12}, 'tickfont': {'color': 'white', 'size': 10}},
            yaxis={'title': 'Average Caudal', 'color': 'white', 'titlefont': {'color': 'white', 'size': 12}, 'tickfont': {'color': 'white', 'size': 10}},
            hovermode='closest',
            plot_bgcolor='#000000',
            paper_bgcolor='#000000',
            margin=dict(l=50, r=0, b=50, t=30),
            height=300,
            width=400,
        )
    }

    return dcc.Graph(
        id='freq-caudal-graph',
        config={'displayModeBar': True},
        style={'display': 'inline-block',
            'backgroundColor': '#000000',
            'color': 'white'},
        figure=figure
    )

# Agrega el título de la tarjeta aquí
presion_caudal_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Intake Pressure vs Flow",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        presion_caudal_graph(df),  # Inserta el mapa directamente después del título
    ],
)

# Agrega el título de la tarjeta aquí
freq_caudal_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Frequency vs Flow",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        freq_caudal_graph(df),  # Inserta el mapa directamente después del título
    ],
)

# Define un diccionario de colores para WCUT y WOR por año
custom_colors_wcut = {
    2021: 'rgba(188, 170, 164, 0.2)', #gris
    2022: 'rgba(0, 204, 255, 0.3)', #azul
    2023: 'rgba(186, 104, 200, 0.2)', #morado
}

custom_colors_wor = {
    2021: 'rgba(255, 255, 157, 0.4)', #amarillo
    2022: 'rgba(255, 87, 34, 0.6)', #naranja
    2023: 'rgba(183, 28, 28, 1)', #rojo
}

# Agrupar los datos por año y mes
df_grouped = df.groupby(['Año', 'Mes']).agg({
    'Suma_WCUT': 'mean',  # Puedes usar 'mean', 'sum', o cualquier otra función de agregación que necesites
    'Suma_WOR': 'mean'
}).reset_index()

# Obtener la lista de años en tus datos
años = df_grouped['Año'].unique()

meses_abreviados = {
        1: 'Jan',
        2: 'Feb',
        3: 'Mar',
        4: 'Apr',
        5: 'May',
        6: 'Jun',
        7: 'Jul',
        8: 'Aug',
        9: 'Sep',
        10: 'Oct',
        11: 'Nov',
        12: 'Dec'
}


# Función para crear gráficas de WCUT/WOR
def create_wc_wor_graph(df_grouped, custom_colors_wcut, custom_colors_wor):
    años = df_grouped['Año'].unique()
    traces = []

    # Actualizar el formato del eje X
    df_grouped['Mes'] = df_grouped['Mes'].map(meses_abreviados) + '/' + df_grouped['Año'].apply(str).str[-2:]


    for año in años:
        data_año = df_grouped[df_grouped['Año'] == año]
        trace_wcut = go.Scatter(
            x=data_año['Mes'],
            y=data_año['Suma_WCUT'],
            mode='lines',
            name=f'WCUT {año}',
            line=dict(
                width=2,
                color=custom_colors_wcut.get(año, '#000033')  # Puedes personalizar los colores aquí
            )
        )
        trace_wor = go.Scatter(
            x=data_año['Mes'],
            y=data_año['Suma_WOR'],
            mode='lines',
            name=f'WOR {año}',
            line=dict(
                width=2,
                color=custom_colors_wor.get(año, '#3366CC')  # Puedes personalizar los colores aquí
            )
        )
        traces.append(trace_wcut)
        traces.append(trace_wor)

    wc_wor_graph = dcc.Graph(
        id='wc-wor-graph',
        config={'displayModeBar': True},
        style={
            'display': 'inline-block',
            'backgroundColor': '#000000',
            'width': '400px',
            'height': '300px',
            'color': 'white',
            'align-items': 'center',
        },
        figure={
            'data': traces,
            'layout': go.Layout(
                xaxis={'title': 'Date',
                        'titlefont': {'size': 12, 'color': 'white'},  # Establece el color del título del eje X en blanco
                        'tickfont': {'size': 9}, 'color': 'white',
                        'showline': True, 'linewidth': 1},  # Muestra la línea del eje X y establece su grosor
                yaxis={'title': 'Units', 'color': 'white', 'titlefont': {'size': 12}, 'tickfont': {'size': 9}},
                hovermode='closest',
                plot_bgcolor='#000000',
                paper_bgcolor='#000000',
                margin=dict(l=50, r=0, b=70, t=30),
                height=300,
                legend=dict(
                    font=dict(size=9, color='white'),  # Configura la fuente de la leyenda
                ),
            )
        }
    )

    return wc_wor_graph

# Define un diccionario de colores para BOPD y BWPD por año
colors_bopd = {
    2021: 'rgba(188, 170, 164, 0.2)', #gris
    2022: 'rgba(0, 204, 255, 0.3)', #azul
    2023: 'rgba(186, 104, 200, 0.2)', #morado
}

colors_bwpd = {
    2021: 'rgba(255, 255, 157, 0.4)', #amarillo
    2022: 'rgba(255, 87, 34, 0.6)', #naranja
    2023: 'rgba(183, 28, 28, 1)', #rojo
}


def create_bopd_bwpd_graph(df, colors_bopd, colors_bwpd):
    años = df['Año'].unique()
    traces_bopd_bwpd = []

    # Actualizar el formato del eje X
    df['Mes'] = df['Mes'].map(meses_abreviados) + '/' + df['Año'].apply(str).str[-2:]

    for año in años:
        data_año = df[df['Año'] == año]
        trace_bopd = go.Scatter(
            x=data_año['Mes'],
            y=data_año['Suma_BOPD'],
            mode='lines',
            name=f'BOPD {año}',
            line=dict(
                width=2,
                color=colors_bopd.get(año, '#9966CC')  # Usar el color del diccionario, o un color predeterminado
            )
        )
        trace_bwpd = go.Scatter(
            x=data_año['Mes'],
            y=data_año['Suma_BWPD'],
            mode='lines',
            name=f'BWPD {año}',
            line=dict(
                width=2,
                color=colors_bwpd.get(año, '#009966')  # Usar el color del diccionario, o un color predeterminado
            )
        )
        traces_bopd_bwpd.append(trace_bopd)
        traces_bopd_bwpd.append(trace_bwpd)

    bopd_bwpd_graph = dcc.Graph(
        id='bopd-bwpd-graph',
        config={'displayModeBar': True},
        style={
            'display': 'inline-block',
            'backgroundColor': '#000000',
            'width': '400px',
            'height': '300px',
            'color': 'white',
            'align-items': 'center',
        },
        figure={
            'data': traces_bopd_bwpd,
            'layout': go.Layout(
                xaxis={'title': 'Date',
                        'titlefont': {'size': 12, 'color': 'white'},  # Establece el color del título del eje X en blanco
                        'tickfont': {'size': 9, 'color': 'white'},
                        'showline': True, 'linewidth': 1},  # Muestra la línea del eje X y establece su grosor
                yaxis={'title': 'BFD', 'color': 'white', 'titlefont': {'size': 12}, 'tickfont': {'size': 9}},
                hovermode='closest',
                plot_bgcolor='#000000',
                paper_bgcolor='#000000',
                margin=dict(l=50, r=0, b=70, t=30),
                height=300,
                legend=dict(
                    font=dict(size=9, color='white'),  # Configura la fuente de la leyenda
                ),
            )
        }
    )

    return bopd_bwpd_graph

# Agrega el título de la tarjeta aquí
wc_wor_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "WC and WOR",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        create_wc_wor_graph(df_grouped, custom_colors_wcut, custom_colors_wor),  # Pasar los argumentos aquí
    ],
)

# Agrega el título de la tarjeta aquí
bopd_bwpd_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "BOPD and BWPD",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        create_bopd_bwpd_graph(df, colors_bopd, colors_bwpd),  # Pasar los argumentos aquí
    ],
)


# Diccionario de colores personalizados por UWI
colors_UWIS = {
    'Well013': 'rgba(188, 170, 164, 0.1)', # Color Gris
    'Well017': 'rgba(67, 160, 71, 0.1)', # Color Verde
    'Well002': 'rgba(255, 87, 34, 0.6)', # Color Naranja
    'Well008': 'rgba(186, 104, 200, 0.6)', # Color Morado
    'Well010': 'rgba(0, 204, 255, 0.6)', # Color Azul
}


# Agregar una columna 'Año_Mes' formateada para ordenar correctamente
well_production_bopd['Año_Mes'] = well_production_bopd['Mes'].apply(lambda x: calendar.month_abbr[x]) + '/' + well_production_bopd['Año'].astype(str).str[-2:]

# Ordenar el DataFrame por 'Año' y 'Mes' de manera descendente
well_production_bopd = well_production_bopd.sort_values(by=['Año', 'Mes'])

def crear_grafico_lineas(well_production_bopd):
    fig = go.Figure()

    for uwi, data in well_production_bopd.groupby('UWI'):
        fig.add_trace(go.Scatter(x=data['Año_Mes'], y=data['Suma_BOPD'], mode='lines', name=uwi, fill='tozeroy', 
                                 line=dict(color=colors_UWIS.get(uwi, 'rgba(255, 255, 255, 0.6)'))))

    fig.update_layout(
        xaxis=dict(title='Month/Year', tickfont=dict(color='white', size=10), titlefont=dict(color='white', size=12),
                   gridcolor='#131313'),
        yaxis=dict(title='Sum Oil Production (BOPD)', tickfont=dict(color='white', size=10), titlefont=dict(color='white', size=12),
                   gridcolor='#131313', dtick=1000000),  # Establece el dtick en 1M
        legend=dict(font=dict(size=10, color='white')),
        width=400, height=300,
        plot_bgcolor='#000000 ',
        paper_bgcolor='#000000 ',
        margin={'t': 50, 'b': 0, 'l': 0, 'r': 0},
        hovermode='closest',
    )

    return fig


def format_month_year(year_month):
    year, month = map(int, year_month.split('-'))
    month_name = calendar.month_abbr[month]
    return f"{month_name}/{str(year)[-2:]}"

# Función para formatear los valores en K (miles) o M (millones)
def format_yaxis(y):
    if y >= 1e6:
        return f'{y / 1e6:.0f}M'
    elif y >= 1e3:
        return f'{y / 1e3:.0f}K'
    else:
        return str(int(y))  # Mostrar valores enteros para valores menores a 1,000

def gas_production(well_production_data):
    # Filtrar los datos para incluir solo los pozos especificados
    wells_to_include = ['Well007', 'Well009', 'Well012', 'Well014', 'Well013']
    filtered_data = well_production_data[well_production_data['UWI'].isin(wells_to_include)]
    well_production_data_sorted = filtered_data.sort_values(by=['UWI', 'Mes'])
    well_production_data_sorted['YearMonth'] = well_production_data_sorted['Año'].astype(str) + '-' + well_production_data_sorted['Mes'].astype(str)
    gas_data = well_production_data_sorted[['UWI', 'YearMonth', 'Suma_Gas', 'Well_Id']]

    custom_colors_gas = {
        'D0A25FCC-4989-4D49-86C1-CAA92F1B3001': 'rgba(67, 160, 71, 0.6)',  # Color verde
        'D0A25FCC-4989-4D49-86C1-CDA92F1B3002': 'rgba(188, 170, 164, 0.6)',  # Color gris
        'D0A25FCC-4989-4D49-86C1-CDA92F1B3004': 'rgba(0, 204, 255, 0.2)',  # Color azul
        'D0A25FCC-4989-4D49-86C1-CDA92F1B3005': 'rgba(183, 28, 28, 0.4)',  # Color rojo
        'D0A25FCC-4989-4D49-86C1-CKA92F1B3006': 'rgba(255, 87, 34, 0.6)',  # Color naranja
        'D0A25FCC-4989-4D49-86C1-CPA92F1B3007': 'rgba(186, 104, 200, 0.2)',  # Color morado
        'D0A25FCC-4989-4D49-86C1-CDA92F1B3008': 'rgba(255, 255, 157, 0.5)',  # Color amarillo
        'D0A25FCC-4989-4D49-86C1-CMN92F1B3009': 'rgba(0, 204, 255, 0.2)',  # Color azul 'rgba(204, 153, 255, 0.2)', # color lila
        '54E7CE87-3AC7-49B2-B794-5730BE7C97010': 'rgba(102, 255, 204, 0.6)', # color celeste
        '971F0184-A90B-4029-99E4-F81C5FAB82012': 'rgba(204, 255, 204, 0.2)', # ccolor menta
        'F95DACDF-1568-4F15-95BF-DE04D3D26013': 'rgba(255, 87, 34, 0.6)',  # Color naranja 'rgba(255, 255, 204, 0.2)', # color piel claro
        'D0A25FCC-4989-4D49-86C1-JKT92F1B3014': 'rgba(67, 160, 71, 0.6)', # color  verde
        'D0A25FCC-4989-4D49-86C1-CDF92F1B3015': 'rgba(255, 153, 204, 0.6)', # color rosado
        'D0A25FCC-4989-4D49-86C1-CAA92F1B3016': 'rgba(153, 204, 102, 0.4)', # color guayabo
        'D0A25FCC-4989-4D49-86C1-CDA92F1B3017': 'rgba(153, 153, 204, 1)', # color azul morado
        'D0A25FCC-4989-4D49-86C1-CDA92F1B3018': 'rgba(97, 97, 97, 0.6)', # color gris oscuro
        'D0A25FCC-4989-4D49-86C1-CDY92F1B3019': 'rgba(141, 110, 99, 0.6)' # color marrón claro
    }

    traces_gas = []

    for well_id, data in gas_data.groupby('Well_Id'):
        uwi = data['UWI'].iloc[0]

        trace_gas = go.Scatter(
            x=data['YearMonth'],
            y=data['Suma_Gas'],
            mode='lines+markers',
            line=dict(color=custom_colors_gas.get(well_id, 'rgba(0, 0, 0, 1)')),
            name=str(uwi),
            fill='tozeroy',
            fillcolor=custom_colors_gas.get(well_id, 'rgba(0, 0, 1, 1)')
        )
        traces_gas.append(trace_gas)

    all_year_months = sorted(well_production_data_sorted['YearMonth'].unique())

    xaxis_labels = [format_month_year(year_month) for year_month in all_year_months]

    max_gas_production = well_production_data_sorted['Suma_Gas'].max()
    max_range_gas = 1000 * ((int(max_gas_production) // 1000) + 1)

    y_tickvals = [0]
    y_ticktext = ['0']
    for i in range(1000, max_range_gas + 1, 1000):
        y_tickvals.append(i)
        y_ticktext.append(f'{i / 1000:.0f}K')

    layout_gas = go.Layout(
        xaxis={'title': 'Month/Year', 'color': 'white', 'tickvals': all_year_months, 'ticktext': xaxis_labels,
               'titlefont': {'size': 12},
               'tickfont': {'size': 9},
        },
        yaxis={'title': 'Sum Gas Production (KPC)', 'color': 'white',
               'titlefont': {'size': 12},
               'tickfont': {'size': 9},
               'tickvals': y_tickvals,
               'ticktext': y_ticktext,
               'dtick': 1000,
               'range': [0, max_range_gas]
        },
        hovermode='closest',
        plot_bgcolor='#000000',
        paper_bgcolor='#000000',
        legend=dict(
            font=dict(size=10, color='white'),
        ),
        margin=dict(l=50, r=50, b=70, t=30),  # Eliminar márgenes
        
    ) 

    gas_graph = dcc.Graph(
        id='gas-production-graph',
        config={'displayModeBar': True},
        style={
            'display': 'inline-block',
            'backgroundColor': '#000000',
            'width': '400px',
            'height': '300px',
            'color': 'white',
            'align-items': 'center',
        },
        figure={
            'data': traces_gas,
            'layout': layout_gas
        }
    )

    return gas_graph

# Define el componente dcc.Graph para mostrar el gráfico
oil_production_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Top 5 highest Oil production",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        dcc.Graph(figure=crear_grafico_lineas(well_production_bopd)),  # Usar dcc.Graph para mostrar el gráfico
    ],
)

# Llama a gas_production() con well_production_data como argumento
gas_production_card = html.Div(
    style={'height': '350px', 'width': '100%', 'textAlign': 'center', 'border-radius': '10px', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Top 5 highest Gas production",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        gas_production(well_production_data),  # Pasa well_production_data como argumento
    ],
)


# Función para formatear números en notación K (miles) y M (millones)
def format_number(number):
    if abs(number) >= 1e6:
        return f"{number / 1e6:.2f} M"
    elif abs(number) >= 1e3:
        return f"{number / 1e3:.2f} K"
    else:
        return f"{number:.2f}"

# Tarjetas de datos
wells_card = generate_data_card("Producing Wells", f"{numero_de_pozos}", '#00CCFF')
# total_oil_card = generate_data_card("Top 5 Oil decline", format_number(df_var_3['Oil'].sum()), '#FF6633')
# total_water_card = generate_data_card("Top 5 Gas reduction", format_number(df_var_3['Water'].sum()), '#0066FF')
# total_gas_card = generate_data_card("Total Gas", format_number(df_var_3['Gas'].sum()), '#663366')
# total_hours_card = generate_data_card("Total Hours", format_number(df_var_3['Hours'].sum()), '#009966')
# average_caudal_card = generate_data_card("Average Caudal", format_number(df_var_2['Caudal'].mean()), ' #CC9900')

total_gas = df_var_3['Gas'].sum()  
total_oil = df_var_3['Oil'].sum()  
total_water = df_var_3['Water'].sum()

# Calcula las proporciones
proporcion_gas = total_gas / (total_gas + total_oil + total_water) * 100
proporcion_oil = total_oil / (total_gas + total_oil + total_water) * 100
proporcion_water = total_water / (total_gas + total_oil + total_water) * 100

# Crea un gráfico de torta de proporción de todos los campos

pie_card = html.Div(
    style={'height': '230px', 'textAlign': 'center', 'border-radius':'7px', 'box-shadow': '0px 0px 5px 2px rgba(255, 255, 255, 0.2)', 'backgroundColor': '#131313'},
    children=[
        html.Div(
            "Oil/Gas/Water Proportion",
            style={
                'paddingTop': '10px',
                'textAlign': 'center',
                'align-content': 'center',
                'color': 'white',
                'borderBottom': '1px solid white',
                'paddingBottom': '10px',
                'fontWeight': 'bold',
            },
        ),
        dcc.Graph(
            id='pie-chart',
            config={'displayModeBar': True},
            style={'height': '170px', 'paddingTop': '10px'},  # Ajusta la altura y el espaciado interno
            figure={
                'data': [
                    go.Pie(
                        labels=['Gas', 'Oil', 'Water'],
                        values=[proporcion_gas, proporcion_oil, proporcion_water],
                        hoverinfo='label+percent',
                        textinfo='percent',
                        textfont=dict(color='black', size=8, family='Lato'),  # Personalizar fuente del porcentaje
                        insidetextfont=dict(color='black', size=14, family='Lato'),  # Personalizar fuente de las etiquetas
                        marker=dict(colors=['#663366', '#FF6633']),
                        hole=0.3,
                    )
                ],
                'layout': go.Layout(
                    plot_bgcolor='#131313',
                    paper_bgcolor='#131313',
                    margin={'t': 0, 'b': 0, 'l': 0.5, 'r': 0},  # Ajusta el margen
                ),
            },
        ),
    ],
)

# ---------------------------------------------------------------------------------------------------------------------------------------------------
# SECCION DE BOTONES DE EXPORTAR
# ---------------------------------------------------------------------------------------------------------------------------------------------------

# Define una función para crear y guardar un archivo de PowerPoint con dos slides personalizados
def create_custom_pptx():
    prs = Presentation()

    # Slide 1: Título personalizado en fondo negro
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco

    # Configurar el fondo negro
    background1 = slide1.background
    fill1 = background1.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0, 0, 0)  # Negro

    # Agregar la imagen de fondo (slide 1)
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    pic1 = slide1.shapes.add_picture("D:/Tablero_prod/assets/spe_logo.png", left, top, width, height)

    # Agregar un cuadro de texto para el título (slide 1)
    left = Inches(1)
    top = Inches(0.1)  # Mover hacia arriba
    width = prs.slide_width - Inches(2)
    height = Inches(1)
    txBox1 = slide1.shapes.add_textbox(left, top, width, height)
    tf1 = txBox1.text_frame
    p1 = tf1.add_paragraph()
    p1.text = "Geohallitians Production Kpis Report"
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.runs[0]
    font1 = run1.font
    font1.size = Pt(36)  # Tamaño de fuente (ajusta según sea necesario)
    font1.bold = True
    font1.color.rgb = RGBColor(255, 255, 255)  # Blanco

    # Captura el área del dashboard usando pyautogui
    screenshot = pyautogui.screenshot()
    screenshot.save("dashboard.png")

    # Slide 2: Título y imagen
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco

    # Agregar la imagen (slide 2)
    left = Inches(0)
    top = Inches(1.5)
    width = prs.slide_width
    height = Inches(5.5)  # Reducir la altura de la imagen
    pic2 = slide2.shapes.add_picture("dashboard.png", left, top, width, height)

    # Agregar el título (slide 2)
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(1)
    txBox2 = slide2.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "KPIs Dashboard"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    font2 = run2.font
    font2.size = Pt(48)  # Aumentar el tamaño de la fuente
    font2.bold = True
    font2.color.rgb = RGBColor(0, 0, 0)  # Negro

    prs.save("Geohallitians_dashboard_presentation.pptx")


# Define una función para capturar el área del dashboard y exportarla a PDF
def export_dashboard_to_pdf():
    # Captura el área del dashboard usando pyautogui
    screenshot = pyautogui.screenshot()
    screenshot.save("dashboard.png")
    
    # Convierte la imagen capturada a PDF usando img2pdf
    with open("Geohallitians_dashboard.pdf", "wb") as f:
        f.write(img2pdf.convert("dashboard.png"))

# Define una función para capturar el área del dashboard y exportarla a PNG
def export_dashboard_to_png():
    # Captura el área del dashboard usando pyautogui
    screenshot = pyautogui.screenshot()
    screenshot.save("Geohallitians_dashboard.png")

# ---------------------------------------------------------------------------------------------------------------------------------------------------
# SECCION DEL LAYOUT DE LA APP
# ---------------------------------------------------------------------------------------------------------------------------------------------------

app.layout = html.Div(style={'fontFamily': 'Lato', 'display': 'flex', 'flexDirection': 'column'}, children=[
    dcc.Location(id='url', refresh=False),  # Agrega el componente dcc.Location
    html.Div(style={'position': 'relative', 'className': 'button-container'}, children=[
        html.Button('Export PDF', id='btn-pdf', n_clicks=0, className='export-button'),
        html.Button('Export PPT', id='btn-pptx', n_clicks=0, className='export-button'),
        html.Button('Export PNG', id='btn-png', n_clicks=0, className='export-button') 
    ]),

    html.Div(style={'display': 'flex', 'backgroundColor' : '#131313', 'border-radius':'7px'}, children=[
        html.Img(src='assets/spe_logo_dim.png', style={'width': '5%', 'align-self': 'flex-start', 'margin-top': '0px', 'margin-bottom': '0px', 'margin-left': '42%'}),
        html.H1("PRODUCTION KPI's", style={'color': 'white',
                'alignItems': 'center',
                'fontWeight': 'bold',
                'flex-grow': '1',
                'fontSize': '30px',
                'margin-left': '2%',
            }),
        html.Div(f'Last Updated: {formatted_time} (Colombia Time)', style={'align': 'center', 'color': '#424242', 'margin': '1px 0', 'align-items': 'top', 'fontSize': '14px', 'fontWeight': 'bold'}),
    ]),

    html.Div(style={'display': 'flex'}, children=[

        #Primera columna

        # Barra de control
        # Primera columna
        html.Div(style={'width': '17%', 'margin-top': '0px'}, children=[
                html.Div(style={'backgroundColor': '#131313', 'border-radius': '7px', 'max-height': '50vh'}, children=[
                    html.H3("Filters", style={'textAlign': 'center', 'color': 'white', 'borderBottom': '1px solid white', 'paddingTop': '10px', 'paddingBottom': '10px'}),
                    # Por ejemplo, puedo agregar un dropdown para seleccionar opciones
                    dcc.Dropdown(
                        id='dropdown-option',
                        options=[
                            # {'label': 'Last 3 months', 'value': 'Last 3 months'},
                            # {'label': 'Last 6 months', 'value': 'Last 6 months'},
                            # {'label': 'Last 9 months', 'value': 'Last 9 months'},
                            # {'label': 'Last Year', 'value': 'Last Year'},
                            {'label': 'All', 'value': 'All'},
                            # Agrego más opciones según mis necesidades más adelante
                        ],
                        value='All',  # Valor por defecto
                        style={'margin-bottom': '10px'}
                    ),
                    # # Otros componentes dcc aquí (botones, sliders, etc.)

                    # dcc.Slider(3, 12, 3, value=12),

                    # # Checkbox 1
                    # dcc.Checklist(
                    #     options=[
                    #         {'label': 'Total', 'value': 'Total'},
                    #         {'label': 'top 5 increase', 'value': 'top 5 increase'},
                    #         {'label': 'top 5 decrease', 'value': 'top 5 decrease'},
                    #     ],
                    #     value=['Total'],
                    #     style={'textAlign': 'center', 'color': 'white', 'margin-top': '15px', 'columnCount': 3, 'margin-bottom': '15px'}
                    # ),

                    # Checkbox 2
                    dcc.Checklist(
                        options=[
                            # {'label': 'Field SPE-2', 'value': 'Field SPE-2'},
                            # {'label': 'Field SPE-3', 'value': 'Field SPE-3'},
                            # {'label': 'Field SPE-4', 'value': 'Field SPE-4'},
                            # {'label': 'Field SPE-5', 'value': 'Field SPE-5'},
                            # {'label': 'Field SPE-6', 'value': 'Field SPE-6'},
                            {'label': 'All Fields', 'value': 'All Fields'},
                        ],
                        value=['All Fields'],
                        style={'textAlign': 'center', 'color': 'white', 'margin-top': '15px', 'columnCount': 2, 'margin-bottom': '15px'},
                    ),

                    # Dropdown pero múltiple
                    dcc.Dropdown(
                        ['BES', 'GAS LIFT', 'BME', 'PCP'], ['BES', 'GAS LIFT', 'BME', 'PCP'], style={'backgroundColor': '#000000', 'margin-top': '5px'},
                        multi=True
                    ),

                    # # Botón final
                    # html.Div([
                    #     dcc.Input(id='input-box', type='text'),
                    #     html.Button('Submit', id='button-example-1'),
                    #     html.Div(id='output-container-button', children='Enter a value and press submit', style={'color': 'white'}),
                    #     html.Div(id='search-result', children='', style={'color': 'white'})
                    # ])
                ]),

                
                # Agregar el gráfico de torta dentro de la misma columna pero separado
                html.Div(style={'margin-top': '30px'}, children=[
                    pie_card,
                ]),

                # Agregar el gráfico de torta dentro de la misma columna pero separado
                html.Div(style={'margin-top': '30px'}, children=[
                    hours_card,
                ]),
                                
                html.Div(style={'margin-top': '30px'}, children=[
                    Barrel_price_realtime_card,
                ]),               
        ]),

        #Segunda columna
        html.Div(style={'margin-left': '30px', 'display': 'flex', 'align-content': 'center', 'flex-direction': 'column', 'margin-top': '20px'}, children=[
            # Primera fila de gráficas
            html.Div(style={'display': 'flex'}, children=[
                wells_location_card,  # Llama a la función para crear el mapa
                html.Div(style={'margin-left': '30px'}, children=[
                    oil_production_card,
                ]),
                html.Div(style={'margin-left': '30px'}, children=[
                    gas_production_card,
                ]),
            ]),

            # Segunda fila de gráficas
            html.Div(style={'display': 'flex', 'margin-top': '30px'}, children=[
                html.Div(style={'width': '50%', 'align': 'center'}, children=[
                    wells_depth_card 
                ]),
                html.Div(style={'margin-left': '30px'}, children=[
                    efficiency_card,
                ]),
                html.Div(style={'margin-left': '30px'}, children=[
                    runstatus_card,
                ]),
            ]),

            # Tercera fila de gráficas
            html.Div(style={'display': 'flex', 'margin-top': '30px'}, children=[
                html.Div(children=[
                    Pruebas_card,
                ]),
                html.Div(style={'margin-left': '30px'}, children=[
                    bopd_bwpd_card,
                ]),
                html.Div(style={'margin-left': '30px'}, children=[
                    API_card,  
                ]),
            ]),

            # Cuarta fila de gráficas
            html.Div(style={'display': 'flex', 'margin-top': '30px'}, children=[
                html.Div(children=[
                    wc_wor_card,
                ]),
                html.Div(style={'margin-left': '30px'}, children=[
                    freq_caudal_card,
                ]),
                html.Div(style={'margin-left': '30px'}, children=[
                     presion_caudal_card, 
                ]),
            ]),

            # Quinta fila de gráficas
            html.Div(style={'display': 'flex', 'margin-top': '30px'}, children=[
                html.Div(children=[
                    heatmap_card
                ]),
            ]),
        ]),

        # Tercera columna
        html.Div(
            style={'width': '13%', 'margin-top': '20px', 'margin-left': '30px'},
            children=[
                wells_card,
                top5minusoil_card,
                top5minusgas_card,
                alertas_card,
                # average_caudal_card
            ],
        ),
    ]),
    html.Footer(children='Julieth Muñoz, Yulitza Parada, Silvio Pacheco (Geohallitians 2023)©', style={'textAlign': 'center', 'color': '#424242', 'align-items': 'end', 'margin': '0'})
])

# @app.callback(
#     Output('search-result', 'children'),
#     Input('button-example-1', 'n_clicks'),
#     Input('input-box', 'value')
# )
# def search_and_highlight(n_clicks, search_text):
#     if n_clicks and search_text:
#         # Obtiene el contenido de la página
#         page_content = dcc._js_dist[0]['namespace']['document']['body']['textContent']

#         # Realiza la búsqueda y obtiene las ubicaciones de todas las coincidencias
#         matches = [m.start() for m in re.finditer(re.escape(search_text), page_content, re.IGNORECASE)]

#         if matches:
#             # Resalta todas las coincidencias en el texto
#             highlighted_text = []
#             start = 0
#             for match in matches:
#                 highlighted_text.append(page_content[start:match])
#                 highlighted_text.append(html.Mark(page_content[match:match+len(search_text)], style={'background-color': 'yellow'}))
#                 start = match + len(search_text)
#             highlighted_text.append(page_content[start:])

#             return highlighted_text

#     return ''

# Callback para manejar los clics en los botones de exportación
@app.callback(
    Output('btn-pdf', 'n_clicks'),
    Output('btn-pptx', 'n_clicks'),
    Output('btn-png', 'n_clicks'),
    [Input('btn-pdf', 'n_clicks'),
     Input('btn-pptx', 'n_clicks'),
     Input('btn-png', 'n_clicks')]
)
def export_dashboard(n_clicks_pdf, n_clicks_pptx, n_clicks_png):
    print(f'n_clicks_pdf: {n_clicks_pdf}')
    print(f'n_clicks_pptx: {n_clicks_pptx}')
    print(f'n_clicks_png: {n_clicks_png}')
    ctx = dash.callback_context

    if not ctx.triggered:
        return [0, 0, 0]

    # Si se hace clic en el botón de exportar a PDF
    if ctx.triggered[0]['prop_id'] == 'btn-pdf.n_clicks':
        export_dashboard_to_pdf()
        return [0, 0, 0]
    
    # Si se hace clic en el botón de exportar a PowerPoint
    if ctx.triggered[0]['prop_id'] == 'btn-pptx.n_clicks':
        create_custom_pptx()
        return [0, 0, 0]

    # Si se hace clic en el botón de exportar a PNG
    elif ctx.triggered[0]['prop_id'] == 'btn-png.n_clicks':
        export_dashboard_to_png()
        return [0, 0, 0]

    return [0, 0, 0]

# Agrega el código JavaScript para el desplazamiento de la página
app.clientside_callback(
    """
    function(url) {
        window.onscroll = function() {
            var exportButtons = document.querySelectorAll('.export-button');
            var scrollTop = window.pageYOffset || document.documentElement.scrollTop;
            
            exportButtons.forEach(function(button, index) {
                // Establece la posición y los márgenes de los botones
                button.style.position = 'absolute';
                button.style.top = scrollTop + 'px';
                button.style.left = (index * 100) + 'px'; // Ajusta el espaciado horizontal
                // Establece el margen derecho solo en el último botón
                if (index === exportButtons.length - 1) {
                    button.style.marginRight = '0px';
                } else {
                    button.style.marginRight = '0'; // Elimina el margen derecho en otros botones
                }
            });
        };
        return null;
    }
    """,
    Output('url', 'pathname'),
    Input('url', 'pathname')
)


if __name__ == '__main__':
    app.run_server(debug=True)
